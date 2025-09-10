# forecast_suite_with_report.py
# -*- coding: utf-8 -*-
"""
Pipeline hoàn chỉnh (Q1 -> Q4) + xuất báo cáo PDF & PowerPoint + đồ họa tinh chỉnh.

Yêu cầu thư viện:
pip install pandas numpy matplotlib scikit-learn statsmodels pyreadstat python-pptx pillow fpdf

Chạy:
python forecast_suite_with_report.py --sav "Du bao bang mo hinh nhan qua san luong _ CPQC Hoahong quy.sav" --h 4 --outdir outputs

Kết quả:
- Thư mục outputs/ chứa biểu đồ PNG tinh chỉnh, CSV dự báo, TXT tóm tắt, file report.pdf và report.pptx
"""
import os
import argparse
import json
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.ticker import MaxNLocator
from PIL import Image

import statsmodels.api as sm
from statsmodels.tsa.holtwinters import SimpleExpSmoothing, Holt, ExponentialSmoothing

# report libs
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

# ------------------ Utility helpers ------------------

def safe_mape(y_true, y_pred):
    y_true = np.asarray(y_true, dtype=float)
    y_pred = np.asarray(y_pred, dtype=float)
    mask = y_true != 0
    if mask.sum() == 0:
        return np.nan
    return np.mean(np.abs((y_true[mask] - y_pred[mask]) / y_true[mask]))

def detect_column(df, candidates, required=False):
    lower_map = {c.lower(): c for c in df.columns}
    for cand in candidates:
        if cand.lower() in lower_map:
            return lower_map[cand.lower()]
    if required:
        raise KeyError(f"Missing required column; tried: {candidates}")
    return None

def ensure_trend_and_quarter(df):
    t_col = detect_column(df, ["t", "trend", "time_index"], required=False)
    if t_col is None:
        df = df.copy()
        df["t"] = np.arange(1, len(df) + 1)
        t_col = "t"
    q_col = detect_column(df, ["quarter", "quy", "q", "Quarter", "Quy"], required=False)
    if q_col is None:
        df = df.copy()
        df["quarter"] = ((df[t_col].astype(int) - 1) % 4) + 1
        q_col = "quarter"
    for q in [2,3,4]:
        df[f"Q{q}"] = (df[q_col] == q).astype(int)
    return df, t_col, q_col

def fit_best_ets(y, seasonal_periods=4):
    results = {}
    y = pd.Series(y).astype(float)
    try:
        ses = SimpleExpSmoothing(y, initialization_method="heuristic").fit()
        results["SES"] = {"model": ses, "aic": getattr(ses, "aic", np.nan)}
    except Exception:
        results["SES"] = {"model": None, "aic": np.inf}
    try:
        holt = Holt(y, initialization_method="heuristic").fit()
        results["Holt"] = {"model": holt, "aic": getattr(holt, "aic", np.nan)}
    except Exception:
        results["Holt"] = {"model": None, "aic": np.inf}
    try:
        hw = ExponentialSmoothing(y, trend="add", seasonal="add", seasonal_periods=seasonal_periods,
                                  initialization_method="heuristic").fit(optimized=True)
        results["HW_add"] = {"model": hw, "aic": getattr(hw, "aic", np.nan)}
    except Exception:
        results["HW_add"] = {"model": None, "aic": np.inf}
    best_name, best_info = min(results.items(), key=lambda kv: kv[1]["aic"])
    return best_name, best_info["model"], results

def train_test_split_time(df, test_h=4):
    if len(df) <= test_h:
        return df.copy(), df.iloc[0:0].copy()
    return df.iloc[:-test_h].copy(), df.iloc[-test_h:].copy()

# ------------------ Plot styling helpers ------------------

# Palette & fonts (you can change these)
PALETTE = {
    "actual": "#1f77b4",     # blue
    "fitted": "#ff7f0e",     # orange
    "forecast": "#2ca02c",   # green
    "grid": "#e6e6e6"
}
TITLE_FONT = {"fontsize": 16, "fontweight": "bold"}
AXIS_FONT = {"fontsize": 12}
TICK_FONT_SIZE = 10
plt.rcParams['font.family'] = 'DejaVu Sans'  # cross-platform default; change if needed

def plot_series_with_forecast(dates, actual, fitted=None, forecast=None, forecast_index=None, title="", ylabel="", outpath=None):
    plt.figure(figsize=(10,4.5))
    ax = plt.gca()
    ax.plot(dates, actual, label="Actual", color=PALETTE["actual"], linewidth=2)
    if fitted is not None:
        ax.plot(dates, fitted, label="Fitted", color=PALETTE["fitted"], linewidth=1.8, linestyle='--')
    if forecast is not None and forecast_index is not None:
        ax.plot(forecast_index, forecast, label="Forecast", color=PALETTE["forecast"], linewidth=2, marker='o')
    ax.set_title(title, **TITLE_FONT)
    ax.set_ylabel(ylabel, **AXIS_FONT)
    ax.grid(axis='y', color=PALETTE["grid"], linestyle='-', linewidth=0.8)
    ax.legend()
    ax.xaxis.set_major_locator(MaxNLocator(integer=True))
    plt.xticks(rotation=20, fontsize=TICK_FONT_SIZE)
    plt.yticks(fontsize=TICK_FONT_SIZE)
    plt.tight_layout()
    if outpath:
        plt.savefig(outpath, dpi=180)
        plt.close()
    else:
        plt.show()

# ------------------ Read .sav ------------------

def read_sav_any(sav_path):
    # Force exact file name as requested by user
    expected_name = "Du bao bang mo hinh nhan qua san luong _ CPQC Hoahong quy.sav"
    base = os.path.basename(sav_path)
    if base != expected_name:
        raise FileNotFoundError(f"File name must be exactly: '{expected_name}'. You provided: '{base}'")
    # Try to read
    try:
        # pandas.read_spss requires pyreadstat
        df = pd.read_spss(sav_path)
        return df
    except Exception:
        try:
            import pyreadstat
            df, meta = pyreadstat.read_sav(sav_path)
            return df
        except Exception as e:
            raise RuntimeError("Không thể đọc file .sav. Cài đặt pyreadstat rồi thử lại.\n" + str(e))

# ------------------ Reporting helpers ------------------

def make_pdf_report(outdir, summary_text, image_paths, metrics_csv):
    pdf_path = os.path.join(outdir, "report.pdf")
    pdf = FPDF(unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=12)
    pdf.add_page()
    pdf.set_font("Arial", size=16, style="B")
    pdf.cell(0, 8, "Báo cáo Dự báo Sản lượng", ln=True, align="C")
    pdf.ln(4)
    pdf.set_font("Arial", size=11)
    # summary
    pdf.multi_cell(0, 6, summary_text)
    pdf.ln(4)
    # images
    for img in image_paths:
        if img and os.path.exists(img):
            pdf.add_page()
            # center image
            pdf.image(img, x=15, w=180)  # scale to width 180mm
    # attach metrics table as text
    pdf.add_page()
    pdf.set_font("Arial", size=12, style="B")
    pdf.cell(0, 8, "Metrics Summary (MAPE, Adj_R2, AIC)", ln=True)
    pdf.set_font("Arial", size=10)
    if os.path.exists(metrics_csv):
        dfm = pd.read_csv(metrics_csv)
        pdf.ln(2)
        for i, row in dfm.iterrows():
            line = " | ".join([f"{c}: {row[c]}" for c in dfm.columns])
            pdf.multi_cell(0, 6, line)
    pdf.output(pdf_path)
    return pdf_path

def make_pptx_report(outdir, summary_text, image_paths, metrics_csv):
    prs = Presentation()
    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Báo cáo Dự báo Sản lượng"
    subtitle.text = "Tự động sinh bởi forecast_suite_with_report.py"
    # summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Tóm tắt"
    body = slide.shapes.placeholders[1].text_frame
    for line in summary_text.splitlines():
        p = body.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
    # image slides
    for img in image_paths:
        if img and os.path.exists(img):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # title + content layout
            slide.shapes.title.text = os.path.basename(img)
            left = Inches(0.5)
            top = Inches(1.4)
            slide.shapes.add_picture(img, left, top, width=Inches(9))
    # metrics slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Metrics Summary"
    if os.path.exists(metrics_csv):
        dfm = pd.read_csv(metrics_csv)
        rows, cols = dfm.shape
        # add a table
        x, y, cx, cy = Inches(0.5), Inches(1.2), Inches(9), Inches(0.8 + 0.2*rows)
        table = slide.shapes.add_table(rows+1, cols, x, y, cx, cy).table
        # header
        for j, c in enumerate(dfm.columns):
            table.cell(0, j).text = c
        for i in range(rows):
            for j, c in enumerate(dfm.columns):
                table.cell(i+1, j).text = str(dfm.iloc[i, j])
    pptx_path = os.path.join(outdir, "report.pptx")
    prs.save(pptx_path)
    return pptx_path

# ------------------ Main pipeline ------------------

def main(args):
    outdir = args.outdir
    os.makedirs(outdir, exist_ok=True)

    # 1) Read data (force exact file name)
    df = read_sav_any(args.sav)
    df.columns = [c.strip() for c in df.columns]

    # 2) detect columns
    y_col = detect_column(df, ["Sanluong", "sanluong", "Output", "Y"], required=True)
    cpqc_col = detect_column(df, ["CPQC", "cpqc"], required=False)
    hh_col = detect_column(df, ["Hoahong", "hoahong", "HoaHong"], required=False)

    # 3) ensure t and quarters
    df, t_col, q_col = ensure_trend_and_quarter(df)

    # 4) build model_df
    model_df = df[[y_col, t_col]].copy()
    if cpqc_col is not None:
        model_df["CPQC"] = df[cpqc_col]
    if hh_col is not None:
        model_df["Hoahong"] = df[hh_col]
    for q in [2,3,4]:
        model_df[f"Q{q}"] = df[f"Q{q}"]

    # create datetime index for plotting (quarterly)
    periods = len(model_df)
    idx = pd.period_range(start=pd.Period("2000Q1", freq="Q"), periods=periods).to_timestamp(how="end")
    model_df.index = idx

    # ---------- Q1: Trend regression ----------
    X1 = sm.add_constant(model_df[[t_col]])
    y = model_df[y_col].astype(float)
    ols1 = sm.OLS(y, X1).fit()
    model_df["yhat_q1"] = ols1.predict(X1)
    mape_q1 = safe_mape(y, model_df["yhat_q1"])

    q1_plot = os.path.join(outdir, "q1_trend_regression.png")
    plot_series_with_forecast(model_df.index, y, fitted=model_df["yhat_q1"],
                              title=f"Q1 - Hồi quy xu thế: {y_col}",
                              ylabel=y_col, outpath=q1_plot)

    with open(os.path.join(outdir, "q1_ols_summary.txt"), "w", encoding="utf-8") as f:
        f.write(ols1.summary().as_text())

    # ---------- Q2: ETS for Sanluong ----------
    best_name_y, best_model_y, _ = fit_best_ets(model_df[y_col].values, seasonal_periods=4)
    fitted_y = best_model_y.fittedvalues if hasattr(best_model_y, "fittedvalues") else best_model_y.predict(start=0, end=len(model_df)-1)
    mape_q2_y = safe_mape(model_df[y_col].values, fitted_y)

    h = args.h
    fc_index = pd.period_range(model_df.index[-1].to_period("Q")+1, periods=h, freq="Q").to_timestamp(how="end")
    fc_y = best_model_y.forecast(h)
    q2_sanluong_forecast = pd.Series(fc_y, index=fc_index, name="Sanluong_forecast")
    q2a_plot = os.path.join(outdir, "q2a_ets_sanluong.png")
    plot_series_with_forecast(model_df.index, model_df[y_col], fitted=fitted_y, forecast=fc_y, forecast_index=fc_index,
                              title=f"Q2a - ETS ({best_name_y}) cho {y_col}", ylabel=y_col, outpath=q2a_plot)
    q2_sanluong_forecast.to_csv(os.path.join(outdir, "q2a_sanluong_forecast.csv"), index_label="date")

    # ---------- Q2b: ETS for regressors ----------
    ets_regressors = {}
    q2b_plots = []
    for col in ["Hoahong", "CPQC"]:
        if col in model_df.columns:
            best_name, best_model, _ = fit_best_ets(model_df[col].astype(float).values, seasonal_periods=4)
            fitted = best_model.fittedvalues if hasattr(best_model, "fittedvalues") else best_model.predict(start=0, end=len(model_df)-1)
            mape = safe_mape(model_df[col].values, fitted)
            fc = best_model.forecast(h)
            ets_regressors[col] = {"best": best_name, "model": best_model, "fitted": fitted, "mape": mape, "forecast": fc}
            pth = os.path.join(outdir, f"q2b_ets_{col.lower()}.png")
            plot_series_with_forecast(model_df.index, model_df[col].values, fitted=fitted, forecast=fc, forecast_index=fc_index,
                                      title=f"Q2b - ETS ({best_name}) cho {col}", ylabel=col, outpath=pth)
            q2b_plots.append(pth)

    # ---------- Q3: Multiple regression ----------
    X_cols = [t_col, "CPQC", "Hoahong", "Q2", "Q3", "Q4"]
    X_cols = [c for c in X_cols if c in model_df.columns]
    X3 = sm.add_constant(model_df[X_cols])
    ols3 = sm.OLS(y, X3).fit()
    model_df["yhat_q3_insample"] = ols3.predict(X3)
    mape_q3_in = safe_mape(y, model_df["yhat_q3_insample"])

    # Build future design using ETS forecasts for regressors if present
    future_df = pd.DataFrame(index=fc_index)
    future_df[t_col] = np.arange(model_df[t_col].iloc[-1] + 1, model_df[t_col].iloc[-1] + h + 1)
    future_quarters = (((future_df[t_col].astype(int) - 1) % 4) + 1).values
    for q in [2,3,4]:
        future_df[f"Q{q}"] = (future_quarters == q).astype(int)
    if "CPQC" in X_cols and "CPQC" in ets_regressors:
        future_df["CPQC"] = ets_regressors["CPQC"]["forecast"]
    if "Hoahong" in X_cols and "Hoahong" in ets_regressors:
        future_df["Hoahong"] = ets_regressors["Hoahong"]["forecast"]

    X3_future = sm.add_constant(future_df[X_cols], has_constant="add")
    q3_forecast = ols3.predict(X3_future)
    q3_forecast = pd.Series(q3_forecast, index=fc_index, name="Sanluong_forecast_q3")
    q3_plot = os.path.join(outdir, "q3_multireg_forecast.png")
    plot_series_with_forecast(model_df.index, model_df[y_col], fitted=model_df["yhat_q3_insample"],
                              forecast=q3_forecast.values, forecast_index=fc_index,
                              title="Q3 - Hồi quy đa biến (dự báo bằng ETS regressors)", ylabel=y_col, outpath=q3_plot)
    q3_forecast.to_csv(os.path.join(outdir, "q3_sanluong_forecast.csv"), index_label="date")
    with open(os.path.join(outdir, "q3_ols_summary.txt"), "w", encoding="utf-8") as f:
        f.write(ols3.summary().as_text())

    # ---------- Q4: Compare models using holdout ----------
    train_df, test_df = train_test_split_time(model_df, test_h=4)

    # Model (2) full (with Hoahong)
    X_cols_2 = [t_col, "CPQC", "Hoahong", "Q2", "Q3", "Q4"]
    X_cols_2 = [c for c in X_cols_2 if c in train_df.columns]
    X_train_2 = sm.add_constant(train_df[X_cols_2])
    ols2_train = sm.OLS(train_df[y_col].astype(float), X_train_2).fit()
    if len(test_df) > 0:
        X_test_2 = sm.add_constant(test_df[X_cols_2], has_constant="add")
        yhat2_test = ols2_train.predict(X_test_2)
        mape_ols2_test = safe_mape(test_df[y_col].values, yhat2_test.values)
    else:
        mape_ols2_test = np.nan

    # Model (3) drop Hoahong
    X_cols_3 = [t_col, "CPQC", "Q2", "Q3", "Q4"]
    X_cols_3 = [c for c in X_cols_3 if c in train_df.columns]
    X_train_3 = sm.add_constant(train_df[X_cols_3])
    ols3alt_train = sm.OLS(train_df[y_col].astype(float), X_train_3).fit()
    if len(test_df) > 0:
        X_test_3 = sm.add_constant(test_df[X_cols_3], has_constant="add")
        yhat3_test = ols3alt_train.predict(X_test_3)
        mape_ols3alt_test = safe_mape(test_df[y_col].values, yhat3_test.values)
    else:
        mape_ols3alt_test = np.nan

    with open(os.path.join(outdir, "q4_model2_summary.txt"), "w", encoding="utf-8") as f:
        f.write(ols2_train.summary().as_text())
    with open(os.path.join(outdir, "q4_model3_summary.txt"), "w", encoding="utf-8") as f:
        f.write(ols3alt_train.summary().as_text())

    # metrics summary
    metrics_summary = pd.DataFrame({
        "Model": ["Q1: Trend", "Q2: ETS (Sanluong) in-sample", "Q3: Full multireg in-sample",
                  "Q4-Compare: Model(2) holdout", "Q4-Compare: Model(3) holdout"],
        "MAPE": [safe_mape(y, model_df["yhat_q1"]), mape_q2_y, mape_q3_in, mape_ols2_test, mape_ols3alt_test],
        "Adj_R2": [np.nan, np.nan, ols3.rsquared_adj, ols2_train.rsquared_adj, ols3alt_train.rsquared_adj],
        "AIC": [ols1.aic, np.nan, ols3.aic, ols2_train.aic, ols3alt_train.aic]
    })
    metrics_csv = os.path.join(outdir, "metrics_summary.csv")
    metrics_summary.to_csv(metrics_csv, index=False)

    # Recommendation
    if pd.notna(mape_ols2_test) and pd.notna(mape_ols3alt_test):
        if mape_ols2_test < mape_ols3alt_test:
            rec = "Chọn Model (2) (có Hoahong) vì MAPE holdout thấp hơn."
        elif mape_ols3alt_test < mape_ols2_test:
            rec = "Chọn Model (3) (không có Hoahong) vì MAPE holdout thấp hơn."
        else:
            rec = "Hai mô hình tương đương theo MAPE holdout; cân nhắc AIC/R² và tính đơn giản."
    else:
        rec = "Không đủ dữ liệu holdout để so sánh khách quan; tham khảo AIC/R² và tính đơn giản."

    # ---------- Save main outputs ----------
    with open(os.path.join(outdir, "q1_ols_summary.txt"), "w", encoding="utf-8") as f:
        f.write(ols1.summary().as_text())

    q1_eq = f"{y_col} = {ols1.params['const']:.4f} + {ols1.params[t_col]:.4f} * t + e"

    # build summary text
    summary_text = (
        f"Phương pháp: Q1→Q4 tự động\n\n"
        f"Q1 (Trend) equation:\n{q1_eq}\nMAPE in-sample: {safe_mape(y, model_df['yhat_q1']):.4f}\n\n"
        f"Q2 (ETS) best for {y_col}: {best_name_y}, MAPE in-sample: {mape_q2_y:.4f}\n\n"
        f"Q3 (Multiple regression) in-sample MAPE: {mape_q3_in:.4f}\n\n"
        f"Q4 Recommendation: {rec}\n"
    )

    # images list in order for report
    image_paths = [q1_plot, q2a_plot, q3_plot] + q2b_plots

    # create PDF and PPTX
    pdf_path = make_pdf_report(outdir, summary_text, image_paths, metrics_csv)
    pptx_path = make_pptx_report(outdir, summary_text, image_paths, metrics_csv)

    # print final info
    print("\n=== DONE ===")
    print("Outputs saved in:", os.path.abspath(outdir))
    print("Key files:")
    print("- Q1 plot:", q1_plot)
    print("- Q2 ATS plot:", q2a_plot)
    if q2b_plots:
        print("- Q2b plots:", q2b_plots)
    print("- Q3 plot:", q3_plot)
    print("- Forecast CSVs:", os.path.join(outdir, "q2a_sanluong_forecast.csv"), os.path.join(outdir, "q3_sanluong_forecast.csv"))
    print("- Metrics summary CSV:", metrics_csv)
    print("- PDF report:", pdf_path)
    print("- PPTX report:", pptx_path)
    print("\nRecommendation (Q4):", rec)
    print("\nEquation Q1:", q1_eq)

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--sav", type=str, required=True,
                        help="Path to the .sav data file. MUST be named exactly: 'Du bao bang mo hinh nhan qua san luong _ CPQC Hoahong quy.sav'")
    parser.add_argument("--h", type=int, default=4, help="Forecast horizon (quarters)")
    parser.add_argument("--outdir", type=str, default="outputs", help="Folder to save outputs")
    args = parser.parse_args()
    main(args)
